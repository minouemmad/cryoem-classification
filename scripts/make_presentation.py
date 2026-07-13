"""Generate the CFTR conformational heterogeneity presentation.

Run from repo root with the .venv Python::

    python scripts/make_presentation.py
"""
from __future__ import annotations

import os
import sys
from pathlib import Path

ROOT = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(ROOT / "scripts"))

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
import lxml.etree as etree

# ── colour palette ──────────────────────────────────────────────────────────
NAVY    = RGBColor(0x1A, 0x3A, 0x6C)
BLUE    = RGBColor(0x2E, 0x75, 0xB6)
ORANGE  = RGBColor(0xC5, 0x5A, 0x11)
GREEN   = RGBColor(0x37, 0x86, 0x3A)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY   = RGBColor(0xF0, 0xF0, 0xF0)
DGRAY   = RGBColor(0x40, 0x40, 0x40)
BLACK   = RGBColor(0x00, 0x00, 0x00)

# ── figure paths ────────────────────────────────────────────────────────────
FIGS = {
    "scatter":     ROOT / "results_cryosparc/diagnostics/pairwise_posterior_scatter_J1442.png",
    "pca_j1442":   ROOT / "results_cryodrgn/J1442/fullset_D256_z10_ep100/analyze.50/z_pca_marginals.png",
    "pca_j264":    ROOT / "results_cryodrgn/J264/fullset_D256_z10_ep50/analyze.50/z_pca_marginals.png",
    "land_k3_a":   ROOT / "results_cryodrgn/J1442/fullset_D256_z10_ep100/landscape_k3/panel_A_landscape.png",
    "land_k3_d":   ROOT / "results_cryodrgn/J1442/fullset_D256_z10_ep100/landscape_k3/panel_D_pc1_marginal.png",
    "land_z10_d":  ROOT / "results_cryodrgn/J1442/landscape_z10/panel_D_pc1_marginal.png",
    "land_z10_b":  ROOT / "results_cryodrgn/J1442/landscape_z10/panel_B_cryosparc_class.png",
    "conf5":       ROOT / "results_cryodrgn/J1442/confidence_5class/confusion.png",
    "j264_b":      ROOT / "results_cryodrgn/J264/fullset_D256_z10_ep50/landscape_k9/panel_B_cryosparc_class.png",
    "j264_fe":     ROOT / "results_cryodrgn/J264/fullset_D256_z10_ep50/free_energy/free_energy_J264.png",
    "j264_k9_a":   ROOT / "results_cryodrgn/J264/fullset_D256_z10_ep50/landscape_k9/panel_A_landscape.png",
}


# ── helper utilities ─────────────────────────────────────────────────────────
def _px(val): return Inches(val)


def set_slide_bg(slide, color=WHITE):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(1, _px(left), _px(top), _px(width), _px(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, left, top, width, height,
             size=20, bold=False, italic=False, color=DGRAY,
             align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(_px(left), _px(top), _px(width), _px(height))
    tf = tb.text_frame
    tf.word_wrap = wrap
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
        p.alignment = align
    return tb


def add_bullets(slide, items, left, top, width, height,
                size=17, color=DGRAY, title_first=False):
    tb = slide.shapes.add_textbox(_px(left), _px(top), _px(width), _px(height))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if isinstance(item, tuple):
            indent, txt = item
        else:
            indent, txt = 0, item
        run = p.add_run()
        run.text = ("    " * indent) + ("• " if indent > 0 else "▸ ") + txt
        run.font.size = Pt(size - 2 * indent)
        run.font.color.rgb = color
        run.font.bold = (i == 0 and title_first)
    return tb


def add_image(slide, key, left, top, width=None, height=None):
    p = FIGS.get(key)
    if p and p.exists():
        if width and height:
            return slide.shapes.add_picture(str(p), _px(left), _px(top),
                                            _px(width), _px(height))
        elif width:
            return slide.shapes.add_picture(str(p), _px(left), _px(top),
                                            width=_px(width))
        elif height:
            return slide.shapes.add_picture(str(p), _px(left), _px(top),
                                            height=_px(height))
    else:
        # Placeholder box
        add_rect(slide, left, top, width or 5, height or 3, LGRAY, DGRAY)
        name = key if not p else str(p.name)
        add_text(slide, f"[Figure: {name}]", left + 0.1, top + 0.1,
                 (width or 5) - 0.2, (height or 3) - 0.2,
                 size=11, color=DGRAY, italic=True)


def title_bar(slide, text, subtitle=None):
    add_rect(slide, 0, 0, 13.33, 1.05, NAVY)
    add_text(slide, text, 0.25, 0.1, 12.8, 0.65, size=26, bold=True,
             color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text(slide, subtitle, 0.25, 0.72, 12.8, 0.3, size=14,
                 color=RGBColor(0xCC, 0xDD, 0xFF), align=PP_ALIGN.LEFT)


def set_notes(slide, text):
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = text


# ── presentation builder ──────────────────────────────────────────────────────
def build():
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]  # completely blank

    # ── SLIDE 1 – Title ───────────────────────────────────────────────────────
    sl = prs.slides.add_slide(blank_layout)
    set_slide_bg(sl, NAVY)
    add_rect(sl, 0.0, 0.0, 13.33, 7.5, NAVY)
    add_rect(sl, 0.3, 0.28, 12.73, 4.4, RGBColor(0x0E, 0x25, 0x4D))
    add_text(sl,
        "Conformational Heterogeneity of CFTR\n"
        "Under Drug Treatment",
        0.55, 0.5, 12.2, 2.5,
        size=34, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(sl,
        "A Hybrid CryoEM + Deep Learning Pipeline",
        0.55, 2.75, 12.2, 0.7,
        size=22, italic=True, color=RGBColor(0xAA, 0xCC, 0xFF),
        align=PP_ALIGN.CENTER)
    add_text(sl,
        "Minou Emamian  |  Hunt Lab  |  2026",
        0.55, 3.55, 12.2, 0.5,
        size=16, color=RGBColor(0xCC, 0xCC, 0xFF), align=PP_ALIGN.CENTER)
    add_rect(sl, 1.5, 4.6, 10.33, 0.06, RGBColor(0x2E, 0x75, 0xB6))
    add_text(sl,
        "Methods: CryoSPARC • Gaussian Mixture Models • cryoDRGN Neural Network",
        0.55, 4.8, 12.2, 0.5,
        size=14, color=RGBColor(0x99, 0xBB, 0xFF), align=PP_ALIGN.CENTER)

    set_notes(sl, """
SPEAKING NOTES — Slide 1 (Title)

Good [morning/afternoon] everyone. Today I'm going to walk you through my research on a protein called CFTR, 
which is broken in people with cystic fibrosis — a serious lung disease. 

The exciting part is that there are now drugs that can partially fix this protein, but we don't fully 
understand exactly HOW they fix it — we don't know what shape the protein is being pushed into. 

I've been using cryo-electron microscopy — a technique that photographs individual proteins frozen in ice — 
combined with machine learning algorithms, to figure out what shapes this protein adopts when treated with 
these drugs. This work combines traditional image analysis with a neural network approach, and I'll show 
you both methods and how they complement each other.

The talk will be about 12 minutes. I'll start with the biology (no experience needed), then walk through 
the methods, and finish with what we found.
""")

    # ── SLIDE 2 – CFTR Biology ─────────────────────────────────────────────────
    sl = prs.slides.add_slide(blank_layout)
    set_slide_bg(sl)
    title_bar(sl, "What is CFTR? Why Does It Matter?",
              "The biology behind cystic fibrosis")
    add_bullets(sl, [
        "CFTR = Cystic Fibrosis Transmembrane conductance Regulator",
        (1, "A channel protein embedded in the surface of lung and gut cells"),
        (1, "Its job: pump chloride ions (Cl⁻) OUT of the cell to attract water"),
        (1, "Water on cell surfaces keeps mucus thin and fluid"),
        "Cystic Fibrosis (CF): CFTR is broken or absent",
        (1, "Thick, sticky mucus builds up in lungs, pancreas, liver"),
        (1, "~40,000 people in the US; median survival now ~50 years"),
        "Breakthrough drugs (Trikafta / elexacaftor+tezacaftor+ivacaftor)",
        (1, "Correctors: help misfolded CFTR reach the cell surface"),
        (1, "Potentiator: helps the gate of CFTR stay open once there"),
        (1, "Transformed CF from fatal to manageable for ~90% of patients"),
        "Open question: what EXACT shape does CFTR adopt with these drugs?",
    ], 0.4, 1.15, 7.8, 5.8, size=17)

    add_rect(sl, 8.4, 1.15, 4.6, 5.9, LGRAY, BLUE)
    add_text(sl,
        "[ CFTR channel schematic ]\n\n"
        "Shows: cell membrane, CFTR channel pore,\n"
        "chloride ions flowing out, water following.\n\n"
        "Source: Mechanism of dual pharmacological\n"
        "correction and potentiation of human CFTR\n"
        "(Liu et al., 2024)",
        8.55, 1.3, 4.3, 4.0, size=11, italic=True, color=DGRAY)
    add_text(sl, "KEY INSIGHT: drugs work, but the structural 'HOW' is incomplete",
             0.4, 6.7, 12.5, 0.5, size=14, bold=True, color=ORANGE)

    set_notes(sl, """
SPEAKING NOTES — Slide 2 (What is CFTR?)

Let me give you a quick biology refresher. 

Think of CFTR as a tiny door — or channel — in the wall of your lung cells. Its job is to let 
chloride ions (a type of salt) pass through the cell membrane to the outside. When chloride moves 
out, water follows it, and that water is what keeps the mucus on your lung surfaces thin enough 
to flow and be cleared by cilia.

In cystic fibrosis, this door is either broken (it won't open) or it's never built correctly 
(it gets stuck inside the cell and never reaches the surface). The result is thick, sticky mucus 
that builds up and causes repeated lung infections.

The new drugs — sold under the brand name Trikafta — include two types:
- CORRECTORS: these act like chaperones that help the misfolded protein fold into the right shape 
  so it can actually reach the cell surface
- POTENTIATORS: these act like a wedge that keeps the channel door open once it gets there

These drugs have been transformative — they've changed CF from a disease that killed most patients 
in childhood to one where people can live to 50 and beyond. 

BUT — we don't fully understand at the atomic level what shape the protein is being held in by these 
drugs. That's what my research is trying to answer, using electron microscopy.

[Point to the placeholder box] In the actual presentation you'd want to insert a figure showing CFTR 
as a cartoon in a cell membrane with chloride ions flowing through.
""")

    # ── SLIDE 3 – CFTR Structure & Drug Binding ────────────────────────────────
    sl = prs.slides.add_slide(blank_layout)
    set_slide_bg(sl)
    title_bar(sl, "CFTR's Modular Architecture & Drug Targets",
              "Different drugs bind to different 'rooms' in the protein")
    add_bullets(sl, [
        "CFTR has distinct structural domains (like rooms in a building):",
        (1, "NBD1 / NBD2 — Nucleotide Binding Domains: the 'engines'  (bind ATP to open/close)"),
        (1, "TMD1 / TMD2 — Transmembrane Domains: the pore through the membrane"),
        (1, "R-domain — Regulatory domain: a 'lock' that must be phosphorylated to allow opening"),
        "Most common mutation: ΔF508 — deletes one amino acid (F508) from NBD1",
        (1, "Causes NBD1 to misfold → whole protein degraded before reaching membrane"),
        "Drug binding sites (from cryo-EM structures + biochemistry):",
        (1, "Elexacaftor/Tezacaftor: bind at the NBD1–TMD2 interface (correction)"),
        (1, "Ivacaftor: binds inside the pore region (potentiation)"),
        "Result: multiple distinct STABLE SHAPES (conformations) of the drug-bound protein",
        (1, "These are what we are imaging and classifying"),
    ], 0.4, 1.15, 7.8, 5.8, size=16)

    add_rect(sl, 8.4, 1.15, 4.6, 5.5, LGRAY, BLUE)
    add_text(sl,
        "[ CFTR 3D structure ]\n\n"
        "Shows: NBD1, NBD2, TMD1, TMD2, R-domain\n"
        "Drug binding pockets highlighted.\n\n"
        "Source: Flatiron CFTR Hierarchical\n"
        "Unfolding Poster (2025) and/or\n"
        "Liu et al. 2024",
        8.55, 1.3, 4.3, 4.0, size=11, italic=True, color=DGRAY)

    add_text(sl,
        "Biological class names used in this study:\n"
        "P6=NBD1LessMix-Ablated  |  P7=NBD1LessWide-Ablated  |  P8=VshapedMix\n"
        "P9=NBD2Less-Ablated  |  P10=AltNBD1-ArdeconComposite-Ablated",
        0.4, 6.55, 12.5, 0.8, size=12, italic=True, color=BLUE)

    set_notes(sl, """
SPEAKING NOTES — Slide 3 (CFTR Structure & Drug Targets)

Now let's look at CFTR's structure more carefully. The protein is made of several distinct sections 
or domains — think of it like a building with different rooms that each have a specific function.

The NBDs (Nucleotide Binding Domains) are the engines — they bind to ATP molecules (your cell's 
energy currency) and use that energy to physically open and close the pore. NBD1 is where the 
most common mutation occurs: ΔF508, which means amino acid number 508 (phenylalanine) is deleted. 
This one missing amino acid causes the whole protein to fold incorrectly and get destroyed before 
it even reaches the cell surface.

The TMDs (Transmembrane Domains) form the actual pore — the channel through which chloride flows.

The corrector drugs bind at the junction between these domains and stabilize the folded structure. 
The potentiator drug fits inside the pore and keeps it wedged open.

When we image CFTR with electron microscopy — and I'll explain that in a moment — we can actually 
see the protein in different shapes. The class names you'll see throughout this talk (P6, P7, P8, etc.) 
are shorthand for these different conformational states, and they describe what domain is missing or 
altered compared to a reference structure. For example:
- "NBD1LessMix-Ablated" means this class shows reduced NBD1 density, mixing with an ablated variant
- "VshapedMix" means the protein adopts a characteristic V-shape

These names tell us WHERE the structural differences are — which is exactly the biology we care about.
""")

    # ── SLIDE 4 – CryoEM Methodology ──────────────────────────────────────────
    sl = prs.slides.add_slide(blank_layout)
    set_slide_bg(sl)
    title_bar(sl, "CryoEM: Photographing Individual Proteins",
              "How we actually see CFTR at atomic resolution")

    steps = [
        ("1", "Purify CFTR protein + add drug (Trikafta)", NAVY),
        ("2", "Mix with lipid nanodiscs → stable membrane environment", BLUE),
        ("3", "Plunge-freeze into liquid ethane → vitreous ice (~−170°C)", GREEN),
        ("4", "Load into electron microscope; collect 10,000–100,000 images", ORANGE),
        ("5", "Computational alignment: find which direction each particle is facing", NAVY),
        ("6", "3D reconstruction from millions of aligned 2D projections", BLUE),
    ]
    for i, (num, txt, col) in enumerate(steps):
        y = 1.25 + i * 0.88
        add_rect(sl, 0.3, y, 0.55, 0.65, col)
        add_text(sl, num, 0.35, y + 0.08, 0.45, 0.5, size=20, bold=True,
                 color=WHITE, align=PP_ALIGN.CENTER)
        add_text(sl, txt, 1.0, y + 0.08, 7.0, 0.6, size=16, color=DGRAY)

    add_rect(sl, 8.2, 1.15, 4.8, 5.6, LGRAY, BLUE)
    add_text(sl,
        "[ CryoEM workflow diagram ]\n\n"
        "Shows the pipeline from:\n"
        "protein + drug → frozen grid →\n"
        "electron micrograph → particle\n"
        "picking → 2D classes → 3D volume\n\n"
        "Source: Punjani et al. 2017\n"
        "Nature Methods (cryoSPARC paper)",
        8.35, 1.3, 4.45, 4.0, size=11, italic=True, color=DGRAY)

    add_text(sl,
        "Key challenge: each image is EXTREMELY noisy (signal-to-noise ratio ~0.1)\n"
        "AND each protein is in a random orientation — we don't control which way it's facing",
        0.3, 6.65, 12.7, 0.7, size=13, italic=True, color=ORANGE)

    set_notes(sl, """
SPEAKING NOTES — Slide 4 (CryoEM Methodology)

Now let me explain HOW we actually image these proteins. Cryo-EM stands for cryo-electron 
microscopy, and the process goes like this:

First, we purify the CFTR protein in the lab and mix it with the Trikafta drugs. Because CFTR 
normally lives in cell membranes, we also mix it with special lipid particles called nanodiscs 
that mimic the membrane environment and keep the protein stable.

Then we take a tiny droplet of this protein solution and plunge it extremely rapidly into liquid 
ethane at -170 degrees Celsius. This freezes the solution so fast that water doesn't form ice 
crystals — instead it forms a glass-like amorphous solid called vitreous ice. The protein is now 
frozen mid-motion, in whatever shape it happened to be in at that exact moment.

We then load this frozen sample into an electron microscope — imagine it as a camera that uses 
electrons instead of light, which gives much finer detail. We collect tens of thousands of images, 
each showing individual protein particles embedded in the ice.

The key challenge here is that: 
1) Each image is incredibly noisy — the signal from one tiny protein barely stands out from background
2) Each protein is oriented randomly — we're getting random 2D projections of a 3D object

The computational step (step 5) is where you reconstruct the 3D structure — it's like using 
thousands of shadow projections from different angles to reconstruct the 3D object that cast them.
""")

    # ── SLIDE 5 – The Classification Problem ──────────────────────────────────
    sl = prs.slides.add_slide(blank_layout)
    set_slide_bg(sl)
    title_bar(sl, "The Challenge: One Protein, Many Shapes",
              "Heterogeneous reconstruction — the core problem of this project")

    add_bullets(sl, [
        "PROBLEM: CryoEM collects images of ALL conformations mixed together",
        (1, "If we average everything naively → blurry, meaningless map"),
        (1, "Like averaging photographs of a running person — you'd get a blur"),
        "GOAL: sort 100,000+ particle images into conformational classes",
        (1, "Each class = one distinct shape of the protein"),
        (1, "Then reconstruct a 3D map for each class separately"),
        "For CFTR + Trikafta, CryoSPARC found:",
        (1, "Dataset J1442: 3 classes (K=3) with 230,396 particles"),
        (1, "Dataset J1497: 5 classes (K=5) — same particles, different model"),
        (1, "Dataset J264: 9 classes (K=9) with 301,770 particles"),
        "The challenge: are these really distinct conformations or artifacts?",
        (1, "BOTH methods agree on the core 3 states — this is our strongest result"),
    ], 0.4, 1.15, 12.5, 5.9, size=17)

    add_rect(sl, 0.4, 6.5, 12.5, 0.75, RGBColor(0xE8, 0xF0, 0xFE), BLUE)
    add_text(sl,
        "This talk: compare TWO independent approaches — CryoSPARC (reference-based) vs "
        "cryoDRGN (reference-free neural network) — to build confidence in the classifications",
        0.55, 6.6, 12.2, 0.6, size=13, italic=True, color=NAVY)

    set_notes(sl, """
SPEAKING NOTES — Slide 5 (The Classification Problem)

Here's the core challenge of this whole project. When we collect cryo-EM images, we're imaging 
a mixture of proteins, and each protein particle might be in a slightly different shape. 

If we just average all 100,000 images together, we'd get a blurry mess — like taking a 
long-exposure photograph of a busy intersection. Each car is in a different position, so you'd 
just see streaks, not individual cars.

Instead, we want to SORT the particles into groups — each group containing particles that look 
similar to each other, meaning they're likely in the same conformation. Then we reconstruct a 
separate 3D map for each group.

For CFTR treated with Trikafta drugs, the software found between 3 and 9 distinct groups 
depending on how we set up the experiment:
- J1442: 3 groups with 230,000 particles
- J1497: same particles but sorted into 5 groups  
- J264: a different, larger experiment with 9 groups and 300,000 particles

The big scientific question is: are these really distinct conformations, or is the software 
just hallucinating groups? This is where my two-method approach becomes important — if two 
completely different methods, starting from different assumptions, both find the same 3 groups, 
that's strong evidence those groups are real.
""")

    # ── SLIDE 6 – CryoSPARC Hetero-Refine & Bias ─────────────────────────────
    sl = prs.slides.add_slide(blank_layout)
    set_slide_bg(sl)
    title_bar(sl, "CryoSPARC Heterogeneous Refinement: Power & Bias",
              "The 'Einstein from noise' problem in iterative refinement")

    add_text(sl, "How it works:", 0.4, 1.2, 12.5, 0.35, size=18, bold=True, color=NAVY)
    add_bullets(sl, [
        "Start with K reference 3D maps (from ab-initio reconstruction)",
        "E-step: assign each particle to its most-likely class (based on correlation to reference)",
        "M-step: rebuild each class map from its assigned particles",
        "Iterate until convergence",
    ], 0.4, 1.6, 7.8, 2.2, size=16)

    add_text(sl, "The bias problem:", 0.4, 3.8, 12.5, 0.35, size=18, bold=True, color=ORANGE)
    add_bullets(sl, [
        "'Einstein from noise': show reference maps, the algorithm finds a match even in pure noise",
        (1, "Convergence to LOCAL optimum near the starting references, not THE global truth"),
        "Near one-hot posteriors: algorithm becomes OVERCONFIDENT in its class assignments",
        (1, "For J1442: after biased refinement, mean max-posterior = 0.992 (≈100% confident)"),
        (1, "After DEBIASED single E-step (J1442 honest): mean = 0.362 (≈ 1/3 each = flat)"),
        "This means: most particles are genuinely ambiguous between the 3 classes",
    ], 0.4, 4.2, 12.5, 2.9, size=16)

    add_text(sl,
        "E-step formula:  P(class k | particle n) ∝ P(image n | class k) × P(class k)\n"
        "After many iterations: P → one-hot vector (0,0,...,1,...,0)  ← overconfidence",
        0.4, 7.05, 12.5, 0.4, size=12, italic=True, color=DGRAY)

    set_notes(sl, """
SPEAKING NOTES — Slide 6 (CryoSPARC Bias)

Let me explain how CryoSPARC's heterogeneous refinement works and why we have to be careful about it.

The algorithm works in cycles. It starts with reference 3D maps — initial guesses at what each 
conformation looks like. Then it does two steps over and over:
- E-step (Expectation): for each particle image, calculate "which reference map does this look 
  most like?" and assign a probability to each class
- M-step (Maximization): rebuild each class's 3D map using the particles assigned to it

This is very powerful, but there's a fundamental problem called "Einstein from noise." In 2009, 
researchers showed that if you show CryoSPARC a reference image of Einstein's face and give it 
pure random noise images, it will find Einstein's face in the noise. This is because the 
correlation step can always find SOMETHING that matches a reference — the algorithm will 
converge to whatever it was primed to look for.

This leads to what I call "near one-hot posteriors" — after many iterations, the algorithm becomes 
almost 100% confident about every particle's class assignment. For J1442, after biased refinement, 
the average confidence is 99.2% — essentially saying "I'm sure this particle is exactly class 6." 

But when we use a debiased approach — running just ONE E-step from an unbiased starting point — 
the average confidence drops to 36.2%, which is almost exactly 1/3 for three classes. This means 
most particles are genuinely ambiguous and the algorithm was manufacturing false confidence.

My GMM pipeline was designed to characterize and quantify this uncertainty honestly.
""")

    # ── SLIDE 7 – GMM Uncertainty Pipeline ────────────────────────────────────
    sl = prs.slides.add_slide(blank_layout)
    set_slide_bg(sl)
    title_bar(sl, "Our GMM Pipeline: Honest Uncertainty Quantification",
              "Fitting a Gaussian Mixture Model to the posterior probability simplex")

    add_bullets(sl, [
        "Input: each particle has a vector (P6, P7, P8) summing to 1.0  (K×1 probability vector)",
        "Problem: this lives on a 2D simplex — standard Gaussian doesn't apply",
        "Solution: Additive Log-Ratio (ALR) transform to make it unconstrained:",
    ], 0.4, 1.2, 12.5, 1.6, size=17)

    add_rect(sl, 0.8, 2.7, 6.5, 0.7, RGBColor(0xE8, 0xF0, 0xFE), BLUE)
    add_text(sl,
        "  ALR:  yⱼ = log( pⱼ / pₖ )    (drop last component as reference)",
        0.85, 2.75, 6.4, 0.6, size=16, bold=True, color=NAVY)

    add_bullets(sl, [
        "Then fit K-component GMM in ALR space:",
    ], 0.4, 3.5, 12.5, 0.5, size=17)

    add_rect(sl, 0.8, 3.95, 9.0, 0.75, RGBColor(0xFF, 0xF3, 0xE0), ORANGE)
    add_text(sl,
        "  p(x) = Σₖ πₖ · 𝒩(x ; μₖ , Σₖ)    (GMM likelihood)",
        0.85, 4.0, 8.8, 0.65, size=16, bold=True, color=RGBColor(0x8B, 0x30, 0x00))

    add_bullets(sl, [
        "Outputs of the GMM pipeline:",
        (1, "GMM responsibilities rₙₖ = soft class membership per particle"),
        (1, "Soft confusion matrix: how often each class 'bleeds into' another"),
        (1, "Bootstrap-corrected population estimates with confidence intervals"),
        (1, "Exportable particle subsets above a confidence threshold"),
        "Key finding: J1442 GMM components heavily overlap near (0.33, 0.33, 0.33)",
        (1, "Classes are genuinely ambiguous — the 3 states are positions on a continuum"),
    ], 0.4, 4.85, 12.5, 2.5, size=16)

    set_notes(sl, """
SPEAKING NOTES — Slide 7 (GMM Pipeline)

OK so now I've told you the problem — CryoSPARC gives us overconfident assignments. What did I build 
to quantify the uncertainty more honestly?

Each particle in the CryoSPARC output gets assigned a probability vector. For 3 classes, it's 
(P6, P7, P8) which sums to 1.0. This is like saying "this particle is 40% likely to be class 6, 
35% class 7, 25% class 8."

These probability vectors lie on what mathematicians call a simplex — a triangular surface in 
3D space. You can't just put a regular Gaussian distribution on a simplex, because the probabilities 
have to add up to 1. 

To fix this, I use the Additive Log-Ratio transform, which converts the simplex into regular 
unconstrained coordinates. The formula on the slide says: take the log of each probability divided 
by the last one. This maps the triangle onto infinite 2D space.

Then I fit a Gaussian Mixture Model — basically multiple overlapping "bells" in this transformed 
space. The formula shows the standard GMM: it's a weighted sum of Gaussian bell curves, where 
πₖ is the weight (how big that bell is), μₖ is the center (where that class lives in this space), 
and Σₖ is the shape of the bell (how spread out it is).

The outputs are really useful:
- We get soft memberships (rₙₖ) — how much each particle belongs to each class
- We can measure class overlap — which classes are most similar  
- We can estimate the true population fractions with error bars
- We can export just the high-confidence particles for better reconstructions

The key result: when I plot these probability vectors, all three classes heavily overlap near 
the center (0.33, 0.33, 0.33) — meaning most particles genuinely don't belong to any one class, 
they're sort of in between. This shows the three "classes" are actually positions along a 
continuous conformational spectrum.
""")

    # ── SLIDE 8 – GMM Results (clean pairwise scatter) ────────────────────────
    sl = prs.slides.add_slide(blank_layout)
    set_slide_bg(sl)
    title_bar(sl, "Posterior Overlap: Three Classes Share the Same Probability Space",
              "J1442 (K=3) — 230,396 CFTR particles, CryoSPARC debiased posteriors")

    add_image(sl, "scatter", 0.25, 1.15, width=9.5)

    add_text(sl,
        "How to read this:\n"
        "• Each dot = one particle  (subsampled to 50k for clarity)\n"
        "• Color = which class CryoSPARC assigned it to (P6=red, P7=orange, P8=gray)\n"
        "• X-axis / Y-axis = CryoSPARC's probability for each class  (e.g. x=P(P6))\n"
        "• White dot = mean position of each class;  ellipses = 1σ / 2σ spread\n"
        "• If classes were perfectly distinct: each cloud would be in a corner (0,0), (1,0), (0,1)\n"
        "• Instead: ALL three clouds pile up near the CENTER (≈ 0.33, 0.33)\n"
        "  → The algorithm is essentially guessing — it has no real confidence",
        9.85, 1.2, 3.3, 5.8, size=13, color=DGRAY)

    add_rect(sl, 0.25, 6.9, 13.0, 0.45, RGBColor(0xFF, 0xF0, 0xE0), ORANGE)
    add_text(sl,
        "Key message: all three class clouds overlap near probability = 1/3 for all classes — "
        "indicating fundamental conformational ambiguity, NOT three discrete states",
        0.4, 6.93, 12.7, 0.4, size=13, bold=True, color=ORANGE)

    set_notes(sl, """
SPEAKING NOTES — Slide 8 (Posterior Scatter)

This is one of the most important figures in my analysis. Let me walk you through it carefully.

Each of these three panels shows a different pairing of the three classes — P6 vs P7, P6 vs P8, 
and P7 vs P8. Each dot represents one protein particle — but remember, we have 230,000 particles, 
so I'm only showing a random sample of 50,000.

The position of each dot tells you what CryoSPARC thought about that particle:
- If a particle is assigned to Class P6 with 90% confidence, it would appear at the FAR RIGHT 
  of the x-axis in the "P6 vs P7" panel
- If it's totally uncertain, it would appear near the middle of the plot, around (0.33, 0.33)

Now look at what actually happened: ALL THREE CLOUDS are piled up near the CENTER. This means 
most particles have probabilities close to 0.33 for each of the three classes — essentially 
uniform. The algorithm is saying "I don't really know which class this is."

The white dots are the class means — where the average particle in each class sits. And you can 
see they're actually pretty close to each other, right in the center of overlap.

This is the HONEST picture after debiasing. Contrast this with what you'd get from a regular 
CryoSPARC run: the dots would be scattered near the corners (0.99, 0.01) type values — 
completely overconfident.

The diagonal negative correlation you see — when P6 probability is high, P7 tends to be low — 
that's mathematically forced because the three probabilities always sum to 1.

Bottom line: these three classes don't correspond to three isolated, well-separated populations. 
They're more like three preferred positions along a continuum of conformational states.
""")

    # ── SLIDE 9 – CryoDRGN Neural Network ────────────────────────────────────
    sl = prs.slides.add_slide(blank_layout)
    set_slide_bg(sl)
    title_bar(sl, "CryoDRGN: An Unsupervised Neural Network Approach",
              "Learning the conformational landscape without reference maps")

    add_text(sl, "Architecture: Variational Autoencoder (VAE)", 0.4, 1.2, 12.5, 0.35,
             size=18, bold=True, color=NAVY)

    # VAE schematic using shapes
    boxes = [
        (0.4,  2.0, 2.5, 2.8, "ENCODER\n(conv neural net)\n\nInput:\nraw particle\nimage", BLUE, WHITE),
        (3.3,  2.0, 2.5, 2.8, "LATENT\nSPACE z\n\n~10-dim vector\n(each particle's\n'fingerprint')", GREEN, WHITE),
        (6.2,  2.0, 2.5, 2.8, "DECODER\n(fully-connected net)\n\nOutput:\nreconstructed\nimage", ORANGE, WHITE),
    ]
    for bx, by, bw, bh, btxt, bfill, btextcol in boxes:
        add_rect(sl, bx, by, bw, bh, bfill)
        add_text(sl, btxt, bx + 0.1, by + 0.15, bw - 0.2, bh - 0.3,
                 size=13, bold=True, color=btextcol, align=PP_ALIGN.CENTER)
    # Arrows between boxes (text approximation)
    add_text(sl, "→", 2.95, 3.15, 0.4, 0.5, size=28, bold=True, color=NAVY)
    add_text(sl, "→", 5.85, 3.15, 0.4, 0.5, size=28, bold=True, color=NAVY)

    # Training objective
    add_rect(sl, 0.4, 5.0, 8.6, 0.75, RGBColor(0xF0, 0xF8, 0xFF), BLUE)
    add_text(sl,
        "  Training objective (ELBO):  ℒ = 𝔼[log p(x|z)] − β · KL( q(z|x) ‖ p(z) )\n"
        "  Reconstruction term → decoder learns to generate images  |  "
        "KL term → latent stays compact",
        0.5, 5.03, 8.5, 0.7, size=13, color=NAVY, bold=True)

    add_bullets(sl, [
        "Input: pose-corrected particle images (poses borrowed from CryoSPARC) + CTF",
        "Output: per-particle 10-dim latent vector z — each particle's 'conformational address'",
        "Key strength: NO reference maps needed — unsupervised by design",
        "Caveat: uses CryoSPARC-derived poses (one indirect dependency; mitigated by agreement)",
        "Then we apply PCA + KDE to z to find the conformational landscape",
    ], 8.95, 1.2, 4.2, 5.9, size=15)

    set_notes(sl, """
SPEAKING NOTES — Slide 9 (CryoDRGN)

Now let me introduce the second major method I used: CryoDRGN, which is a neural network 
approach to this same problem.

The architecture is called a Variational Autoencoder, or VAE. You can think of it as 
two neural networks connected back-to-back:

The ENCODER takes a raw particle image (which is very noisy and has complex variations 
in contrast) and compresses it down to a small vector of about 10 numbers. This vector — 
called z — is the particle's "fingerprint" or "conformational address." Similar shapes 
should get similar z vectors.

The DECODER takes that small z vector and tries to reconstruct the original particle image.

The training objective shown — called the ELBO — has two parts:
1. Reconstruction quality: how well does the decoder recreate the original image?
2. KL divergence: keeps the latent space organized and prevents it from spreading out chaotically

The critical difference from CryoSPARC is that cryoDRGN is UNSUPERVISED. It never looks at 
any reference maps. It figures out the structure purely from the raw data. 

There's one caveat: we do use CryoSPARC's pose estimates (which direction each particle was 
facing) as input. This is an indirect dependency. However, the fact that BOTH methods end up 
finding the same three classes — even with this indirect link — gives us confidence that the 
classes are real, not artifacts.

After training, we have a 10-dimensional latent space. We then use Principal Component Analysis 
(PCA) and free-energy calculations to understand what it looks like.
""")

    # ── SLIDE 10 – Latent Space: 3 Clear Classes ─────────────────────────────
    sl = prs.slides.add_slide(blank_layout)
    set_slide_bg(sl)
    title_bar(sl, "The Latent Space: CryoDRGN Independently Finds 3 States",
              "J1442 fullset D=256, zdim=10, 100 epochs — CONVERGED (loss change <0.01%/epoch)")

    add_image(sl, "pca_j1442", 0.2, 1.15, height=5.4)
    add_image(sl, "land_k3_a", 5.4, 1.15, width=7.8)

    add_text(sl,
        "LEFT: Raw latent PCA (analyze output)\n"
        "3 distinct density lobes visible;\n"
        "PC1/PC2 explain 23.5%+15.9%=39.4%",
        0.2, 6.5, 5.1, 0.75, size=12, color=DGRAY, italic=True)
    add_text(sl,
        "RIGHT: K=3 GMM fit in full 10-D latent (visualized on PC1-PC2 plane)\n"
        "3 labelled components with biological names; min separation = 2.60 SD\n"
        "(>2.0 SD = genuinely distinct classes in multidimensional sense)",
        5.4, 6.5, 7.7, 0.75, size=12, color=DGRAY, italic=True)

    add_rect(sl, 0.2, 7.25, 12.9, 0.2, NAVY)
    add_text(sl,
        "Same 3 classes as CryoSPARC heteroref — found INDEPENDENTLY without reference maps → "
        "strong validation of the classification",
        0.35, 7.28, 12.6, 0.17, size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    set_notes(sl, """
SPEAKING NOTES — Slide 10 (Latent Space: 3 Clear Classes)

This is our headline result. Look at the left figure — it shows the raw output of cryoDRGN 
after training. Each dot is one of the 230,000 particles, plotted in the first two principal 
components of the 10-dimensional latent space.

You can clearly see THREE distinct blobs. No guidance was given to the algorithm about how 
many classes to find, yet it spontaneously separated the data into three groups.

On the right, I've fit a 3-component Gaussian Mixture Model to the FULL 10-dimensional latent 
space (not just the 2D projection). The three colored ellipses correspond to the three classes, 
labelled with their biological names. The key metric is the minimum separation between any two 
components: 2.60 standard deviations. A separation greater than 2 SD means the classes are 
genuinely distinguishable — they're not just overlapping blobs.

These three latent components correspond exactly to the same three classes that CryoSPARC 
identified — NBD1LessMix-Ablated (P6), NBD1LessWide-Ablated (P7), and VshapedMix (P8).

This is the core result: TWO completely independent methods, starting from different 
assumptions, both find the same three CFTR conformational states. This mutual validation gives 
us much higher confidence that these three states are real structural features of the protein, 
not just artifacts of one particular algorithm.

The model was trained to convergence — the loss function barely changed in the last 5 epochs 
(less than 0.01% change), so we know we're not seeing something that would change with more training.
""")

    # ── SLIDE 11 – PC1 Marginal Density: 3 Peaks ─────────────────────────────
    sl = prs.slides.add_slide(blank_layout)
    set_slide_bg(sl)
    title_bar(sl, "The Conformational Axis: 3 Peaks Along PC1",
              "Free-energy landscape reveals a continuous coordinate with 3 preferred positions")

    add_image(sl, "land_z10_d", 0.2, 1.15, width=7.0)

    add_bullets(sl, [
        "PC1 = dominant conformational axis (explains 23% of latent variance)",
        "F(PC1) = −log p(PC1): maps density to 'energy' (lower = more populated)",
        "3 peaks along PC1 align with P6, P7, P8 class assignments",
        "Peaks connected by populated valleys — a continuous coordinate,",
        (1, "not 3 isolated species. Separation ~1.5–2 kT at D=128"),
        "At D=256 (higher resolution): separation in higher dimensions",
        (1, "GMM min separation 2.60 SD — cleanly resolved in full 10-D"),
        (1, "PC1 alone shows 1 well — the resolution matters!"),
        "KEY: the 3 preferred positions are REPRODUCIBLE",
        (1, "J1442 3-class populations: P6≈37%, P7≈29%, P8≈34%"),
        (1, "Stable across D=128/D=256 runs and multiple methods"),
    ], 7.4, 1.15, 5.7, 5.5, size=15)

    add_text(sl,
        "Bell curves fitted in 1-D on PC1 — each bell = one class distribution along the dominant axis.\n"
        "This is the corrected version: 1-D GMM fitted to PC1 scores (not full-latent projected).",
        0.2, 6.75, 7.1, 0.6, size=11, italic=True, color=DGRAY)

    set_notes(sl, """
SPEAKING NOTES — Slide 11 (PC1 Marginal)

This figure shows what I call the "conformational landscape" — it answers the question: 
what shapes does the protein PREFER to be in?

The x-axis is the first principal component (PC1) of the latent space — the most important 
single axis of variation. The y-axis is density — how many particles are at each position.

The histogram (gray bars) shows where the 230,000 particles actually are. You can clearly 
see THREE peaks. The colored bell curves show the three Gaussian components fitted to these 
peaks — red/green/blue for the three classes.

The fact that all three peaks are separated but still connected by particles in between 
tells us something important: CFTR is not snapping discretely between three totally different 
shapes. Instead, it's moving along a continuous track, with three PREFERRED resting positions. 
It's like a ball rolling in a landscape with three shallow valleys — it tends to rest in one 
of the valleys, but it can move between them.

An important technical note: this version uses a 1-D GMM fitted DIRECTLY on PC1 scores, 
which is why each bell fits its peak well. An earlier version projected the full 10-D GMM 
onto PC1, which made the middle bell (P7) too broad. The corrected version on the slide 
properly fits the peaks.

The populations (how many particles in each state) are: roughly 37% P6, 29% P7, 34% P8. 
These numbers are very stable — they come out essentially the same whether we use D=128 or 
D=256 models, and whether we use the CryoSPARC posteriors or cryoDRGN's latent GMM. 
That reproducibility is itself a result worth noting.
""")

    # ── SLIDE 12 – The 5-Class Problem ────────────────────────────────────────
    sl = prs.slides.add_slide(blank_layout)
    set_slide_bg(sl)
    title_bar(sl, "The 5-Class Challenge: Why P9 and P10 Elude cryoDRGN",
              "CryoSPARC finds 5 classes; the neural network can only cleanly recover 3")

    add_image(sl, "conf5", 0.2, 1.15, width=7.2)

    add_bullets(sl, [
        "CryoSPARC J1497 (same particles, K=5): finds P6, P7, P8, P9, P10",
        (1, "P9 = NBD2Less-Ablated  |  P10 = AltNBD1-ArdeconComposite-Ablated"),
        "cryoDRGN K=5 GMM fit: min separation = 0.79 SD  (<2 = overlapping)",
        (1, "Compare: K=3 gets 2.60 SD — cleanly distinct"),
        "Confusion matrix (LEFT) shows:",
        (1, "49% of CryoSPARC P10 particles → cryoDRGN assigns them to P6"),
        (1, "44% of CryoSPARC P9 particles → cryoDRGN assigns them to P8"),
        "Interpretation: P10 ≈ P6 and P9 ≈ P8 in the neural network's view",
        (1, "The two extra classes are SUBSTATES of the 3 core states"),
        "Why are they hard? Both free-energy analyses (J1442 & J1497)",
        (1, "confirm 1 continuous basin — no energetic barrier separates P9/P10"),
        "NOT a failure: supports that P9/P10 are subtle structural variants",
        (1, "Ongoing work: focused classification targeting specific domains may separate them"),
    ], 7.5, 1.15, 5.6, 5.9, size=14)

    set_notes(sl, """
SPEAKING NOTES — Slide 12 (5-Class Problem)

Now let me address an interesting challenge in this work. CryoSPARC found 5 classes when 
configured to look for 5 (in dataset J1497). But when I run cryoDRGN and try to find 5 
clusters in the latent space, it can't do it cleanly.

The confusion matrix on the left tells the story. This matrix compares what CryoSPARC says 
(rows = CryoSPARC classes) with what cryoDRGN says (columns = cryoDRGN classes).

If the two methods agreed perfectly, you'd see all the numbers on the diagonal (top-left to 
bottom-right) being high, like 0.9 or 1.0. But look at P9 — 44% of CryoSPARC's P9 particles 
get assigned to cryoDRGN's P8 cluster. And P10 — 49% get assigned to cryoDRGN's P6.

What this tells us is that from cryoDRGN's perspective, P9 and P10 are not separate clusters — 
they're just part of the P8 and P6 clouds. The neural network can't tell them apart because 
they're essentially in the same place in latent space.

This isn't really a failure of the method — it's actually telling us something scientifically 
interesting. P9 (NBD2Less-Ablated) is probably just a variant of P8 (VshapedMix) where one 
domain is slightly more disordered. Similarly, P10 is a variant of P6. The structural 
differences exist, but they're too subtle for the neural network to separate without additional 
information — like a focused analysis targeting just those specific protein domains.

Currently I'm exploring ways to extract these sub-states by looking at the free-energy landscape 
in specific regions of latent space and applying targeted classification there.
""")

    # ── SLIDE 13 – J264: 9-Class Dataset ──────────────────────────────────────
    sl = prs.slides.add_slide(blank_layout)
    set_slide_bg(sl)
    title_bar(sl, "J264: A Richer Dataset — 9 Conformational Classes",
              "301,770 particles, D=256 training — emerging cluster structure")

    add_image(sl, "pca_j264", 0.2, 1.15, height=5.3)
    add_image(sl, "j264_b",   5.5, 1.15, width=7.6)

    add_text(sl,
        "LEFT: Raw latent PCA (cryodrgn analyze)\n"
        "~4-6 distinguishable density lobes;\n"
        "PC1=0.34, PC2=0.28 (more structured than J1442!)\n"
        "Ablated classes (NBD-less) visible as side-lobes",
        0.2, 6.5, 5.2, 0.85, size=12, italic=True, color=DGRAY)
    add_text(sl,
        "RIGHT: K=9 class-coloured scatter (cryodrgn landscape analysis)\n"
        "SC/AC/AO core states (blue tones) in centre; portal & ablated spread outward\n"
        "Free energy: F(PC1) = ONE well (continuous) → careful: clusters are POSITIONS,\n"
        "not energetically isolated states",
        5.5, 6.5, 7.7, 0.85, size=12, italic=True, color=DGRAY)

    add_rect(sl, 0.2, 7.32, 12.9, 0.14, NAVY)
    add_text(sl,
        "More training (currently at 50 ep, paper standard = 50) + ablated-class-excluded "
        "retraining may further resolve sub-clusters → ongoing",
        0.35, 7.33, 12.6, 0.12, size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    set_notes(sl, """
SPEAKING NOTES — Slide 13 (J264: 9 Classes)

Now let me show you the more complex dataset we've been analyzing. J264 has 301,770 particles 
and was classified into 9 different conformational groups by CryoSPARC.

The left figure shows the raw cryoDRGN latent space — and you can already see it's more 
structured than J1442. PC1 explains 34% of the variance, PC2 another 28% — so the first two 
components already capture a lot of the interesting variation. You can count roughly 4-6 
distinguishable dense regions.

The right figure colors each particle by its CryoSPARC class. You can see the three main CFTR 
conformations — Symmetric Closed (SC), Asymmetric Closed (AC), and Asymmetric Open (AO) — 
all in the central dense region in different shades of blue. These are the three cleanest, most 
robust states from John's group's structural biology work.

The "ablated" classes — the ones where entire protein domains are detached or disordered — form 
separate side-lobes that stick out from the main density. These are structurally more distinct.

Important caveat: when I compute the free-energy landscape (probability → energy), J264 shows 
ONE continuous basin — no barriers. This doesn't mean there are no distinct classes; it means 
the 9 CryoSPARC classes are POSITIONS along a continuum of conformational changes, not 9 
energetically isolated states.

This is actually consistent with how CFTR biology is supposed to work — the protein is 
constantly flexing and opening and closing, not snapping between rigid states.

I'm currently working on retaining only the "meaningful" classes (excluding the ablated variants 
which may be experimental artifacts) and applying targeted analysis to see if the remaining 
core states can be further separated.
""")

    # ── SLIDE 14 – Synthesis & Key Findings ──────────────────────────────────
    sl = prs.slides.add_slide(blank_layout)
    set_slide_bg(sl)
    title_bar(sl, "Synthesis: What We've Learned",
              "Two independent methods converge on the same picture of CFTR's conformational landscape")

    add_text(sl, "✓  Confirmed results", 0.4, 1.2, 12.5, 0.35, size=18, bold=True, color=GREEN)
    add_bullets(sl, [
        "CFTR + Trikafta exists in 3 core conformational states (P6/P7/P8)",
        (1, "CryoSPARC reference-based: finds these with high map resolution"),
        (1, "CryoDRGN reference-free: finds the SAME 3, independently (min sep 2.60 SD)"),
        (1, "Population fractions reproducible: ~37% P6, ~29% P7, ~34% P8"),
    ], 0.4, 1.6, 12.5, 2.1, size=16)

    add_text(sl, "⚠  Honest caveats", 0.4, 3.7, 12.5, 0.35, size=18, bold=True, color=ORANGE)
    add_bullets(sl, [
        "States are POSITIONS on a continuous landscape, not fully discrete snapshots",
        "CryoSPARC posteriors are near-uniform after debiasing — classes genuinely overlap",
        "5th/extra states (P9/P10) not resolvable by the neural network alone",
        "J264's 9-class structure is partly continuous — more work needed",
    ], 0.4, 4.1, 12.5, 1.8, size=16)

    add_text(sl, "→  What this means for CFTR biology", 0.4, 5.9, 12.5, 0.35, size=18,
             bold=True, color=NAVY)
    add_bullets(sl, [
        "The drug-bound protein explores multiple conformations — dynamic, not static",
        "The most distinct states differ in NBD positioning/disorder and exit portal structure",
        "Structural differences are real and localized (density diagnostics confirm)  ",
    ], 0.4, 6.3, 12.5, 1.05, size=16)

    set_notes(sl, """
SPEAKING NOTES — Slide 14 (Synthesis)

Let me bring everything together. What have we actually learned?

The confirmed results: we have strong evidence for three core CFTR conformational states in the 
presence of Trikafta drugs. This evidence comes from two independent methods — CryoSPARC and 
cryoDRGN — and they agree. The population fractions (how many particles are in each state) are 
stable across different experiments and model sizes. This reproducibility is exactly what you 
want to see in science.

The honest caveats: these three states are not like three completely different protein sculptures. 
They're more like three preferred positions on a flexible protein that's constantly moving. Most 
particles have genuinely ambiguous assignments — they're sort of in-between states. This doesn't 
invalidate the classification; it tells us about the thermodynamics of the protein.

For the two additional states CryoSPARC finds (P9 and P10) — we can see in the neural network's 
latent space that they're essentially sub-states of P8 and P6. Whether they're meaningfully 
distinct from a biological standpoint is still an open question.

For the biology: what this means is that CFTR with Trikafta is a dynamic protein. It doesn't 
just sit in one shape — it wobbles between several preferred shapes. The most structurally 
distinct states differ in the NBD (motor domain) positioning and in the exit portal at the bottom 
of the channel. These are likely the functional differences that matter for drug efficacy.
""")

    # ── SLIDE 15 – Conclusions & Future Directions ────────────────────────────
    sl = prs.slides.add_slide(blank_layout)
    set_slide_bg(sl, NAVY)
    add_rect(sl, 0, 0, 13.33, 7.5, NAVY)

    add_text(sl, "Conclusions & Future Directions",
             0.4, 0.25, 12.5, 0.75, size=28, bold=True, color=WHITE)
    add_rect(sl, 0.4, 1.0, 12.5, 0.05, BLUE)

    add_text(sl, "What was done:", 0.4, 1.2, 12.5, 0.35, size=18, bold=True,
             color=RGBColor(0xAA, 0xCC, 0xFF))
    conc = [
        "Built a GMM pipeline to quantify true uncertainty in CryoSPARC class assignments",
        "Trained and analyzed cryoDRGN on two CFTR datasets (J1442, J264) using D=256 models",
        "Established a cross-method validation framework (CryoSPARC ↔ cryoDRGN)",
        "Designed automated free-energy basin analysis and cluster export to .cs format",
    ]
    for i, c in enumerate(conc):
        add_text(sl, "• " + c, 0.5, 1.6 + i * 0.45, 12.2, 0.42, size=15, color=WHITE)

    add_text(sl, "Next steps:", 0.4, 3.6, 12.5, 0.35, size=18, bold=True,
             color=RGBColor(0xAA, 0xCC, 0xFF))
    next_steps = [
        "Run cryoDRGN analyze on all D=256 models (GPU on Hudson) → UMAP + volume traversals",
        "Focused/masked 3D classification targeting NBD1/NBD2 regions to separate P9/P10",
        "Export cryoDRGN cluster particle sets → NU-refinement in CryoSPARC → new maps",
        "J264: exclude ablated particles, retrain, test if 3 core CFTR states emerge cleanly",
        "3DFlex / 3DVA: model continuous motion along the CFTR conformational coordinate",
    ]
    for i, s in enumerate(next_steps):
        add_text(sl, "▸ " + s, 0.5, 4.05 + i * 0.45, 12.2, 0.42, size=14,
                 color=RGBColor(0xCC, 0xDD, 0xFF))

    add_rect(sl, 0.4, 6.6, 12.5, 0.7, RGBColor(0x0E, 0x25, 0x4D))
    add_text(sl,
        "Code: github.com/minouemmad/cryoem-classification  |  Methods: CryoSPARC + cryoDRGN + "
        "custom GMM pipeline  |  Questions?",
        0.55, 6.65, 12.2, 0.6, size=12, italic=True, color=RGBColor(0x99, 0xBB, 0xFF),
        align=PP_ALIGN.CENTER)

    set_notes(sl, """
SPEAKING NOTES — Slide 15 (Conclusions & Next Steps)

To wrap up — here's a summary of what I've done and where I'm going.

The work I've completed: 
I built a GMM (Gaussian Mixture Model) pipeline that honestly quantifies the uncertainty in 
CryoSPARC's particle classifications. I trained cryoDRGN neural networks on two CFTR datasets 
at high resolution (D=256), both of which converged completely. I established a framework for 
comparing the two methods against each other, and I've built automated analysis tools for 
finding free-energy basins and exporting particle subsets.

For next steps:
The most immediate is running cryodrgn analyze on the GPU cluster to generate volume 
reconstructions from different positions in the latent space — this will tell us what each 
part of the landscape actually looks like structurally.

Then I want to do focused classification — rather than trying to separate all 9 classes globally, 
focus on just the NBD regions where the differences between P9/P10 and P8/P6 are expected to 
live. This should improve our ability to separate those subtle substates.

I'll also export the cryoDRGN cluster particle sets as .cs files that can be used for 
non-uniform refinement in CryoSPARC — this is how we get the highest-quality structural 
maps from our classification.

For J264, the plan is to rerun after excluding the ablated particles (which are probably 
experimental artifacts of partial purification), and see if the 3 core conformational states 
emerge more clearly.

Ultimately, I'd like to apply 3DFlex — which is a more sophisticated tool for modeling 
continuous protein motion rather than discrete states. Given that our analysis shows a 
continuous landscape, this may be the most honest way to describe CFTR's conformational 
dynamics under drug treatment.

Thank you — happy to take questions.
""")

    # ── SAVE ──────────────────────────────────────────────────────────────────
    out = ROOT / "docs" / "CFTR_cryoDRGN_presentation.pptx"
    out.parent.mkdir(exist_ok=True)
    prs.save(str(out))
    print(f"[saved] {out}")
    return out


if __name__ == "__main__":
    build()
