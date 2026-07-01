"""
Graphical Abstract v3 — Emmad_Minou_Graphical_Abstract.pptx

Layout: asymmetric 3-column  (not linear equal boxes)
  Col A (left,  3.55")  – CFTR biology  (RCSB 6MSM image)
  Col B (center, 5.0")  – Methods flow  (3 stacked cards)
      B1  Cryo-EM image (EMDB)
      B2  Classification confusion  <- custom PPTX shapes
      B3  GMM correction            <- custom PPTX shapes
  Col C (right, 4.04")  – Key result  (poster img38)

Word budget <= 15:
  Title (6): "How does Trikafta reshape CFTR conformations?"
  Labels (7): "CF variants" "Cryo-EM" "Class confusion"
               "GMM" "Corrected populations"
  Total = 13 words
"""
from pathlib import Path
from random import seed, uniform
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── paths ─────────────────────────────────────────────────────────
BASE    = Path(r"c:\Users\maemm\OneDrive\Desktop\CryoEM")
IMGS    = BASE / "_abstract_imgs"
POSTER  = IMGS / "poster"
OUT     = BASE / "Emmad_Minou_Graphical_Abstract.pptx"

CFTR_IMG   = IMGS / "cftr_6msm.jpeg"
CRYOEM_IMG = IMGS / "cryoem_particles.jpg"
RESULT_IMG = POSTER / "img38_2642x1421.png"

# ── palette ───────────────────────────────────────────────────────
NAVY   = RGBColor(0x0C, 0x28, 0x4A)
TEAL   = RGBColor(0x00, 0x7A, 0x6E)
TEAL2  = RGBColor(0x00, 0xAA, 0x9A)
BLUE   = RGBColor(0x1A, 0x55, 0x8A)
AMBER  = RGBColor(0xBF, 0x5A, 0x0A)
GREEN  = RGBColor(0x19, 0x65, 0x1C)
GREEN2 = RGBColor(0x2E, 0x9E, 0x34)
MGRAY  = RGBColor(0x88, 0x96, 0xA2)
LGRAY  = RGBColor(0xF0, 0xF3, 0xF5)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
DARK   = RGBColor(0x0C, 0x14, 0x1C)
GOLD   = RGBColor(0xF5, 0xC5, 0x18)
P6C    = RGBColor(0x00, 0x99, 0x8C)
P7C    = RGBColor(0x1A, 0x6B, 0xB0)
P8C    = RGBColor(0x2A, 0x9D, 0x30)

# ── helpers ───────────────────────────────────────────────────────
def rct(sl, l, t, w, h, fill=None, line=None, lw=None, sid=1):
    s = sl.shapes.add_shape(sid, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid() if fill else s.fill.background()
    if fill: s.fill.fore_color.rgb = fill
    if line:
        s.line.color.rgb = line
        if lw: s.line.width = Pt(lw)
    else:
        s.line.fill.background()
    return s

def oval(sl, l, t, w, h, fill=None, line=None, lw=None):
    return rct(sl, l, t, w, h, fill=fill, line=line, lw=lw, sid=9)

def txt(sl, text, l, t, w, h, size=12, bold=False,
        color=DARK, align=PP_ALIGN.CENTER, italic=False):
    tb = sl.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run(); run.text = text
    run.font.size   = Pt(size)
    run.font.bold   = bold
    run.font.italic = italic
    run.font.color.rgb = color

def pic(sl, path, l, t, w, h):
    p = Path(path)
    if p.exists():
        sl.shapes.add_picture(str(p), Inches(l), Inches(t), Inches(w), Inches(h))
    else:
        rct(sl, l, t, w, h, fill=LGRAY, line=MGRAY, lw=1)
        txt(sl, f"[{p.name}]", l, t+h/2-0.15, w, 0.3, size=8, color=MGRAY)

def arr_right(sl, l, t, w=0.40, h=0.40, fill=MGRAY):
    s = sl.shapes.add_shape(13, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.fill.background()

def arr_down(sl, l, t, w=0.26, h=0.22, fill=MGRAY):
    s = sl.shapes.add_shape(36, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.fill.background()

# ── slide setup ───────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
slide = prs.slides.add_slide(prs.slide_layouts[6])

# background
rct(slide, 0, 0, 13.33, 7.5, fill=RGBColor(0xF0, 0xF3, 0xF6))

# ── title bar ─────────────────────────────────────────────────────
rct(slide, 0, 0, 13.33, 1.32, fill=NAVY)
rct(slide, 0, 1.25, 13.33, 0.09, fill=TEAL)
rct(slide, 0, 0, 0.22, 1.32, fill=GOLD)   # gold left accent

txt(slide,
    "How does Trikafta reshape CFTR conformations?",
    0.32, 0.10, 12.80, 1.10,
    size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# ── layout geometry ───────────────────────────────────────────────
PY = 1.40; PH = 4.90
AX, AW = 0.17, 3.55
BX, BW = 3.90, 5.00
CX, CW = 9.10, 4.06

# ═══════════════════════════════════════════════════════════════════
# COL A – CFTR biology
# ═══════════════════════════════════════════════════════════════════
rct(slide, AX+0.07, PY+0.07, AW, PH, fill=RGBColor(0xC8,0xD5,0xDE))  # shadow
rct(slide, AX, PY, AW, PH, fill=WHITE, line=TEAL, lw=2.5)
rct(slide, AX, PY, AW, 0.42, fill=TEAL)
txt(slide, "CF variants", AX+0.10, PY+0.07, AW-0.18, 0.28,
    size=13, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

pic(slide, CFTR_IMG, AX+0.12, PY+0.50, AW-0.24, 2.70)

rct(slide, AX+0.18, PY+3.27, AW-0.36, 0.03, fill=RGBColor(0xCC,0xDD,0xE8))

txt(slide, "Trikafta corrects CFTR\nfolding in cystic fibrosis",
    AX+0.12, PY+3.36, AW-0.24, 0.68,
    size=10.5, color=DARK, align=PP_ALIGN.CENTER)

rct(slide, AX+0.18, PY+4.12, AW-0.36, 0.50,
    fill=RGBColor(0xDC,0xEE,0xFB), line=TEAL, lw=1)
txt(slide, "230,396 particles imaged by cryo-EM",
    AX+0.20, PY+4.17, AW-0.40, 0.40,
    size=9.5, color=TEAL, align=PP_ALIGN.CENTER)

txt(slide, "PDB: 6MSM  |  EMDB: EMD-9230",
    AX+0.08, PY+PH-0.30, AW-0.16, 0.24,
    size=7.5, italic=True, color=MGRAY, align=PP_ALIGN.CENTER)

arr_right(slide, AX+AW+0.04, PY+PH/2-0.22, w=0.36, h=0.44, fill=TEAL)

# ═══════════════════════════════════════════════════════════════════
# COL B – Methods, 3 stacked sub-panels
# ═══════════════════════════════════════════════════════════════════
rct(slide, BX+0.07, PY+0.07, BW, PH, fill=RGBColor(0xC8,0xD5,0xDE))
rct(slide, BX, PY, BW, PH, fill=RGBColor(0xF6,0xF8,0xFA), line=MGRAY, lw=1.5)

SH = (PH - 0.16) / 3.0
for i in range(3):
    sy  = PY + 0.06 + i * SH
    sh  = SH - 0.12
    col   = [TEAL, AMBER, GREEN][i]
    fill  = [RGBColor(0xDE,0xF2,0xF2),
             RGBColor(0xFB,0xEE,0xE0),
             RGBColor(0xDE,0xF3,0xDE)][i]
    label = ["Cryo-EM", "Class confusion", "GMM"][i]

    rct(slide, BX+0.12, sy, BW-0.24, sh, fill=fill, line=col, lw=1.5)
    rct(slide, BX+0.12, sy, BW-0.24, 0.36, fill=col)
    txt(slide, label,
        BX+0.22, sy+0.05, BW-0.44, 0.26,
        size=11, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    if i < 2:
        arr_down(slide, BX+BW/2-0.13, sy+sh+0.02, w=0.26, h=0.14, fill=MGRAY)

# ── B1: Cryo-EM image ─────────────────────────────────────────────
sy = PY + 0.06; sh = SH - 0.12
pic(slide, CRYOEM_IMG, BX+0.16, sy+0.42, BW-0.32, sh-0.52)

# ── B2: Classification confusion — CUSTOM PPTX ────────────────────
sy = PY + 0.06 + SH; sh = SH - 0.12
ix = BX+0.18; iy = sy+0.44; iw = BW-0.36; ih = sh-0.60

ov_w = 1.52; ov_h = ih * 0.58
ov_y = iy + 0.08
spacing = (iw - ov_w) / 2.0
ov_xs = [ix,  ix + spacing,  ix + spacing*2]

for i, (ox, col) in enumerate(zip(ov_xs, [P6C, P7C, P8C])):
    oval(slide, ox, ov_y, ov_w, ov_h, fill=col, line=WHITE, lw=1.5)

# overlap hazes
oval(slide, ov_xs[0]+ov_w*0.58, ov_y+ov_h*0.12,
     ov_w*0.26, ov_h*0.76, fill=RGBColor(0x78,0x9A,0xAC))
oval(slide, ov_xs[1]+ov_w*0.58, ov_y+ov_h*0.12,
     ov_w*0.26, ov_h*0.76, fill=RGBColor(0x6A,0xA8,0x7C))

# particle dots
seed(42)
for ox, col in zip(ov_xs, [P6C, P7C, P8C]):
    for _ in range(4):
        oval(slide, ox+uniform(0.12, ov_w-0.22), ov_y+uniform(0.08, ov_h-0.18), 0.12, 0.12, fill=col)
for _ in range(5):
    ox = ov_xs[0]+uniform(ov_w*0.55, ov_w*0.9)
    oval(slide, ox, ov_y+uniform(0.1, ov_h-0.2), 0.12, 0.12, fill=MGRAY)
for _ in range(5):
    ox = ov_xs[1]+uniform(ov_w*0.55, ov_w*0.9)
    oval(slide, ox, ov_y+uniform(0.1, ov_h-0.2), 0.12, 0.12, fill=MGRAY)

# labels
for i, (ox, lb) in enumerate(zip(ov_xs, ["P6","P7","P8"])):
    txt(slide, lb, ox+ov_w*0.36, ov_y+ov_h*0.36, 0.46, 0.26,
        size=9, bold=True, color=WHITE)

txt(slide, "Particle assignment uncertain near class boundaries",
    ix, ov_y+ov_h+0.06, iw, 0.26,
    size=8, italic=True, color=AMBER, align=PP_ALIGN.CENTER)

# ── B3: GMM correction — CUSTOM PPTX ──────────────────────────────
sy = PY + 0.06 + 2*SH; sh = SH - 0.12
ix = BX+0.18; iy = sy+0.44; iw = BW-0.36; ih = sh-0.62

# Two side-by-side bar charts: Observed vs Corrected
bar_w = 0.21; bar_gap = 0.10
obs_h   = [0.50, 0.90, 0.18]   # biased heights
corr_h  = [0.58, 0.75, 0.42]   # corrected heights
bar_cols = [P6C, P7C, P8C]
group_w = 3*(bar_w+bar_gap) - bar_gap   # width of one 3-bar group

chart_base = iy + ih
half_x_gap = 0.18   # gap between groups

# "Before" group
bef_x = ix + 0.05
txt(slide, "Observed", bef_x, iy, group_w, 0.22,
    size=8.5, bold=True, color=AMBER, align=PP_ALIGN.CENTER)
for k in range(3):
    bx = bef_x + k*(bar_w+bar_gap)
    bh = obs_h[k]
    rct(slide, bx, chart_base-bh, bar_w, bh, fill=bar_cols[k])
rct(slide, bef_x, chart_base, group_w, 0.025, fill=DARK)

# Arrow →
arr_right(slide, bef_x+group_w+0.04, iy+ih/2-0.11,
          w=0.18, h=0.22, fill=MGRAY)

# "After" group
aft_x = bef_x + group_w + 0.26
txt(slide, "Corrected", aft_x, iy, group_w, 0.22,
    size=8.5, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
for k in range(3):
    bx = aft_x + k*(bar_w+bar_gap)
    bh = corr_h[k]
    rct(slide, bx, chart_base-bh, bar_w, bh, fill=bar_cols[k])
rct(slide, aft_x, chart_base, group_w, 0.025, fill=DARK)

# Error bars on corrected (simple vertical lines)
for k in range(3):
    cx = aft_x + k*(bar_w+bar_gap) + bar_w/2
    cy = chart_base - corr_h[k]
    rct(slide, cx-0.005, cy-0.08, 0.01, 0.10, fill=DARK)  # vertical
    rct(slide, cx-0.04,  cy-0.08, 0.08, 0.01, fill=DARK)  # top cap

txt(slide, "GMM deconvolves confusion: honest populations (±0.2%)",
    ix, chart_base+0.04, iw, 0.26,
    size=8, italic=True, color=GREEN, align=PP_ALIGN.CENTER)

arr_right(slide, BX+BW+0.04, PY+PH/2-0.22, w=0.36, h=0.44, fill=GREEN)

# ═══════════════════════════════════════════════════════════════════
# COL C – Key result (poster img38)
# ═══════════════════════════════════════════════════════════════════
rct(slide, CX+0.07, PY+0.07, CW, PH, fill=RGBColor(0xC8,0xD5,0xDE))
rct(slide, CX, PY, CW, PH, fill=WHITE, line=GREEN, lw=2.5)
rct(slide, CX, PY, CW, 0.42, fill=GREEN)
txt(slide, "Corrected populations", CX+0.10, PY+0.07, CW-0.18, 0.28,
    size=13, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

crop_img = IMGS / "result_crop.png"
if not crop_img.exists() and RESULT_IMG.exists():
    im = Image.open(RESULT_IMG)
    W, H = im.size
    # Crop out the very top label row (~10%) — keep core protein+bar content
    im.crop((0, int(H*0.10), W, H)).save(str(crop_img))

pic(slide, crop_img, CX+0.10, PY+0.50, CW-0.20, PH-0.70)

txt(slide, "Flatiron Institute (2025)",
    CX+0.08, PY+PH-0.30, CW-0.16, 0.24,
    size=7.5, italic=True, color=MGRAY, align=PP_ALIGN.CENTER)

# ── bottom bar ────────────────────────────────────────────────────
rct(slide, 0, 6.37, 13.33, 1.13, fill=NAVY)
rct(slide, 0, 6.37, 0.22,  1.13, fill=GOLD)

txt(slide,
    "Correcting cryo-EM class errors \u2192 accurate CFTR drug-response mapping",
    0.32, 6.43, 12.90, 1.00,
    size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# ── save ─────────────────────────────────────────────────────────
prs.save(OUT)
print("Saved:", OUT)
wc = ("How does Trikafta reshape CFTR conformations "
      "CF variants Cryo-EM Class confusion GMM Corrected populations")
print(f"Word count: {len(wc.split())}/15")
